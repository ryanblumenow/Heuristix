import streamlit as st
from streamlit_option_menu import option_menu
import base64
from email import header
from html.entities import html5
from importlib.resources import read_binary
# import hydralit as hy
from markdown import markdown
from numpy.core.fromnumeric import var
import streamlit
import streamlit as st
import sys
from streamlit.web import cli as stcli
from PIL import Image
from functions import *
import streamlit.components.v1 as components
import pandas as pd
from st_clickable_images import clickable_images
import numpy as np
import statsmodels.api as sm
from numpy import mean
from numpy import std
from sklearn.datasets import make_classification
from sklearn.model_selection import cross_val_score
from sklearn.model_selection import RepeatedStratifiedKFold
from sklearn.linear_model import LogisticRegression
from matplotlib import pyplot
from sklearn.model_selection import train_test_split
import matplotlib.pyplot as plt
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler
from sklearn.metrics import silhouette_score
from scipy.spatial.distance import cdist
import seaborn as sns
from io import BytesIO
from statsmodels.formula.api import ols
# from streamlit.state.session_state import SessionState
# import tkinter
import matplotlib
# matplotlib.use('TkAgg')
# matplotlib.use('Agg')
from sklearn.neural_network import MLPClassifier
from sklearn.metrics import accuracy_score
from mpl_toolkits.mplot3d import Axes3D
from matplotlib import cm
from matplotlib.ticker import LinearLocator, FormatStrFormatter
from sklearn.tree import DecisionTreeRegressor, plot_tree
import sklearn
from sklearn.datasets import make_classification
from sklearn.ensemble import RandomForestClassifier
from sklearn.datasets import make_regression
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import RepeatedKFold
import time
from scipy.stats import pearsonr
from scipy.stats import spearmanr
import dtale
from dtale.views import startup
from dtale.app import get_instance
import webbrowser
import dtale.global_state as global_state
import dtale.app as dtale_app
from matplotlib.pyplot import axis, hist
from scipy import stats as stats
# from bioinfokit.analys import stat
from statsmodels.stats.anova import AnovaRM
import statsmodels.api as sm
from statsmodels.graphics.factorplots import interaction_plot
from sklearn.decomposition import PCA
from st_click_detector import click_detector
from streamlit.components.v1 import html
import streamlit.components.v1 as components
import pandas as pd
import streamlit as st
from click_image_copy_from_demo import st_click_image
from dtale.views import startup
from streamlit_quill import st_quill
# import pandas_profiling
from streamlit_pandas_profiling import st_profile_report
import dataingestion
from sklearn.ensemble import RandomForestClassifier
from sklearn.feature_selection import SelectFromModel
from sklearn.preprocessing import LabelEncoder
from streamlit_extras.add_vertical_space import add_vertical_space

st.set_page_config(layout="wide")

def analytix():

    title = '<p style="font-family:sans-serif; color:gold; font-size: 39px; text-align: center;"><b>Automated Analytix flow</b></p>'
    st.markdown(title, unsafe_allow_html=True)

    if 'edanotes' not in st.session_state:
        st.session_state.edanotes = "Write your notes on EDA here"

    if 'correlnotes' not in st.session_state:
        st.session_state.correlnotes = "Write your notes on correlation analysis here"

    if 'hyptestingnotes' not in st.session_state:
        st.session_state.hyptestingnotes = "Write your notes on hypothesis testing here"

    if 'dimrednotes' not in st.session_state:
        st.session_state.dimrednotes = "Write your notes on dimension reduction here"

    if 'anovarednotes' not in st.session_state:
        st.session_state.anovarednotes = "Write your notes on ANOVA analysis here"

    if 'regnotes' not in st.session_state:
        st.session_state.regnotes = "Write your notes on regression analysis here"

    if 'clusternotes' not in st.session_state:
        st.session_state.clusternotes = "Write your notes on cluster analysis here"

    if 'conjointnotes' not in st.session_state:
        st.session_state.conjointnotes = "Write your notes on conjoint analysis here"

    if 'neuralnetworksrednotes' not in st.session_state:
        st.session_state.neuralnetworksrednotes = "Write your notes on neural networks here"

    if 'dectreenotes' not in st.session_state:
        st.session_state.dectreenotes = "Write your notes on decision tree analysis here"

    ### From Jupyter - 0. Prepare the data

    # df, df2, branddf = dataingestion.readdata()
    df = dataingestion.readdata()
    butcol1_, butcol2_, butcol3_ = st.columns([1,1,9.5])
    with butcol1_:
        viewdatasample = st.button("View sample of data", key="datapreview")
    with butcol2_:
        cleardatasample = st.button("Clear sample view of data", key="cleardatapreview")
    with butcol3_:
      pass

    dfcol1, dfcol2, dfcol3 = st.columns([1,3,1])
    
    if viewdatasample:
        with dfcol2:
            st.markdown("### <center>Sample of data:</center>", unsafe_allow_html=True)
            st.write(df.sample(8))
        add_vertical_space(3)
    # df=df.sample(frac=0.03, replace=True, random_state=1)
    # df2=df2.sample(frac=0.03, replace=True, random_state=1)
    # branddf=branddf.sample(frac=0.03, replace=True, random_state=1)

    if cleardatasample:
        with dfcol2:
            st.write("")

    # vars = list(df2.columns.values.tolist())
    vars = list(df.columns.values.tolist())

    st.session_state['pagechoice'] = 'auto analytics'

    # with st.sidebar:
    #     choose = option_menu("ABI Analytics", ["Keep in mind", "Show me the steps", "Give me tips", "Give feedback"],
    #                         icons=['key', 'bezier2', 'joystick', 'keyboard'],
    #                         menu_icon="app-indicator", default_index=0,
    #                         styles={
    #         "container": {"padding": "5!important", "background-color": "#fafafa"},
    #         "icon": {"color": "black", "font-size": "25px"}, 
    #         "nav-link": {"font-size": "16px", "text-align": "left", "margin":"0px", "--hover-color": "#eee"},
    #         "nav-link-selected": {"background-color": "#ab0202"},
    #     }
    #     )

    #     print(st.session_state.pagechoice)
        
    #     if choose == "Keep in mind":
    #         st.write("Remember to ask good questions. That is the basis of making good decisions.")

    #     if choose == "Show me the steps" and st.session_state.pagechoice=="test":
    #         st.write("Steps you should follow:")

    #     if choose == "Give me tips":
    #         st.write("Here are some tips:")

    #     if choose == "Give feedback":
    #         st.write("Give feedback")
    #         with st.sidebar.form(key='columns_in_form',clear_on_submit=True): #set clear_on_submit=True so that the form will be reset/cleared once it's submitted
    #             rating=st.slider("Please rate the app", min_value=1, max_value=5, value=3,help='Drag the slider to rate the app. This is a 1-5 rating scale where 5 is the highest rating')
    #             text=st.text_input(label='Please leave your feedback here')
    #             submitted = st.form_submit_button('Submit')
    #             if submitted:
    #                 st.write('Thanks for your feedback!')
    #                 st.markdown('Your Rating:')
    #                 st.markdown(rating)
    #                 st.markdown('Your Feedback:')
    #                 st.markdown(text)

    if df is not None:

        value = st_click_image()
        if value is None:
            # st.stop()
            pass

        st.success("{} selected".format(value))

        col1, col2, col3 = st.columns([1,0.6,1])

        clicked = ""

        with col1:
            pass

        # with col2:

        #     content2 = """
        #     <link rel="stylesheet" href="https://www.w3schools.com/w3css/4/w3.css">
        #         <style>
        #         .zoom {
        #         padding: 10px;
        #         width: 290px;
        #         height: 180px;
        #         transition: transform .21s; /* Animation */
        #         margin: 0 auto;
        #         }
        #         .zoom:hover {
        #         transform: scale(1.1); /* (110% zoom - Note: if the zoom is too large, it will go outside of the viewport) */
        #         }
        #         </style>
        #         <div class="w3-container">
        #         <a href='#' id='Get started'><img width='21%' src='https://i.postimg.cc/2yyc69HY/analyticslbl.jpg' class='zoom w3-round-xxlarge w3-hover-opacity'></a>
        #         <!--<div class="w3-display-bottom w3-large"><b><h3>Click to start</h3></b></div>-->
        #         </div>
        #         """

        #     clicked = click_detector(content2)

            # st.markdown(f"**{clicked} selected - begin analysis below.**" if clicked != "" else "         **Select a methodology and click 'Get started'!**")

        with col3:
            pass

        # if clicked == "Get started":

        if value=='EDA':
            
            st.subheader("Exploratory data analysis")

            with st.spinner("Analyzing and summarizing dataset and generating dataset profile"):

                my_bar = st.progress(0)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                edaenv = st.expander("Guidance on EDA", expanded=False)

                with edaenv:

                    st.info("User guide")

                    def show_pdf(file_path):
                        # Opening tutorial from file path
                        with open(file_path, "rb") as f:
                            base64_pdf = base64.b64encode(f.read()).decode('utf-8')

                        # Embedding PDF in HTML
                        pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="1500" height="800" type="application/pdf"></iframe>'

                        # Displaying File
                        st.markdown(pdf_display, unsafe_allow_html=True)

                    col1, col2,col3= st.columns(3)
                    with col1:  
                        if st.button('Read PDF tutorial',key='1'):            
                            show_pdf('.\Automated flow\DtaleInstructions-compressed.pdf')
                    with col2:
                        st.button('Close PDF tutorial',key='2')                   
                    with col3:
                        with open(".\Automated flow\DtaleInstructions-compressed.pdf", "rb") as pdf_file:
                            PDFbyte = pdf_file.read()
                        st.download_button(label="Download PDF tutorial", key='3',
                                data=PDFbyte,
                                file_name="EDA Instructions.pdf",
                                mime='application/octet-stream')

                datadescrip = st.expander("Description of data")

                with datadescrip:

                    st.write(df.describe(include='all'))
                    
                edaprofiling = st.expander("Profile of dataset", expanded=False)
                
                with edaprofiling:
                
                    # @st.cache(allow_output_mutation=True)
                    # def gen_profile_report(df, *report_args, **report_kwargs):
                    #     return df.profile_report(*report_args, **report_kwargs)

                    # pr = gen_profile_report(df, explorative=True)

                    # st_profile_report(pr)

                    # df.drop(columns='BRANDNAME', axis=1)

                    @st.cache(allow_output_mutation=True)
                    def gen_profile_report(df, *report_args, **report_kwargs):
                        return df.profile_report(*report_args, **report_kwargs)

                    pr = gen_profile_report(df, explorative=True, title="Data profile",
                    dataset={
                    "description": "This profiling report shows an overview of the data",
                    "copyright_holder": "Heuristix",
                    "copyright_year": "2024",
                    "url": "https://www.ryanblumenow.com"}, vars={"num": {"low_categorical_threshold": 0}} )

                    st_profile_report(pr)

            # startup(data_id="1", data=df2.sample(15000)) # All records, no OHE
            startup(data_id="1", data=df)

            if get_instance("1") is None:
                startup(data_id="1", data=df)

            d=get_instance("1")

            # webbrowser.open_new_tab('http://localhost:8501/dtale/main/1') # New window/tab
            # components.html("<iframe src='/dtale/main/1' />", width=1000, height=300, scrolling=True) # Element
            html = f"""<iframe src="/dtale/main/1" height="1000" width="1400"></iframe>""" # Iframe
            # html = "<a href='/dtale/main/1' target='_blank'>Dataframe 1</a>" # New tab link

            st.markdown(html, unsafe_allow_html=True)

            # d = dtale.show(pd.DataFrame(df2.sample(1000)))
            st.session_state.corr_img = d.get_corr_matrix()
            st.session_state.corr_df = d.get_corr_matrix(as_df=True)
            st.session_state.pps_img = d.get_pps_matrix()
            st.session_state.pps_df = d.get_pps_matrix(as_df=True)

            print(st.session_state.corr_df)

            checkbtn = st.button("Validate data")

            if checkbtn == True:
                df_amended = get_instance(data_id="1").data # The amended dataframe ready for upsert
                st.write("Sample of amended data:")
                st.write("")
                st.write(df_amended.head(5))

            clearchanges = st.button("Clear changes made to data")
            if clearchanges == True:
                global_state.cleanup()

            st.write("")
            
            st.subheader("Notes on EDA")

            # Spawn a new Quill editor
            st.subheader("Notes on exploratory data analysis")
            edacontent = st_quill(placeholder="Write your notes here", value=st.session_state.edanotes, key="edaquill")

            st.session_state.edanotes = edacontent

            st.write("Exploratory data analysis took ", time.time() - start_time, "seconds to run")

        elif value == "Correlation analysis":

            st.subheader("Correlation analysis")

            with st.spinner('Please wait while we conduct the correlation analysis'):

                my_bar = st.progress(0)

                time.sleep(10)

                for percent_complete in range(100):
                    time.sleep(0.1)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                x=0
                y=0

                x = st.selectbox("Please choose first variable:", vars, key="correlx", index=x)
                y = st.selectbox("Please choose second variable:", vars, key="correly", index=y)

                # var1 = df[x].iloc[1:101]
                # var2 = df[y].iloc[1:101]

                var1 = df[x].sample(10000)
                var2 = df[y].sample(10000)

                st.subheader("Correlation between chosen 2 variables:" + " " + x + " and " + y)

                st.write("Correlation coefficient: ", var1.corr(var2))
                st.write("Pearson correlation coefficient: ", pearsonr(var1, var2))
                st.write("Spearman correlation coefficient: ", spearmanr(var1, var2))

                ### From Jupyter: 1. Hypothesis testing: correlation

                st.subheader("Correlation matrix, Pearson")

                st.write(df.corr(method = 'pearson'))

                with st.expander("Interpretation guide"):
                    st.write("Variables within a dataset can be related for many reasons.\n\n"
                    "For example:\n"
                    "One variable could cause or depend on the values of another variable.\n"
                    "One variable could be lightly associated with another variable.\n"
                    "Two variables could depend on a third unknown variable.\n\n"
                    "It can be useful in data analysis and modeling to better understand the relationships between variables. The statistical relationship between two variables is referred to as their correlation.\n"
                    "A correlation could be positive, meaning both variables move in the same direction, or negative, meaning that when one variable's value increases, the other variables' values decrease. Correlation can also be neutral or zero, meaning that the variables are unrelated.")
                    st.info("Positive Correlation: both variables change in the same direction.\n\n"
                    "Neutral Correlation: No relationship in the change of the variables.\n\n"
                    "Negative Correlation: variables change in opposite directions.\n\n"
                    "A correlation coefficient closer to 1 or to -1 indicates stronger relationships between the variables, while a correlation coefficient closer to 0 indicates a lesser relationship between the variables.")
                    st.write("The performance of some algorithms can deteriorate if two or more variables are tightly related, called multicollinearity. An example is linear regression, where one of the offending correlated variables should be removed in order to improve the skill of the model.\n\n"
                    "We may also be interested in the correlation between input variables with the output variable in order provide insight into which variables may or may not be relevant as input for developing a model.\n"
                    "The structure of the relationship may be known, e.g. it may be linear, or we may have no idea whether a relationship exists between two variables or what structure it may take. Depending what is known about the relationship and the distribution of the variables, different correlation scores can be calculated.")

                ### End

            st.write("")

            # Spawn a new Quill editor
            st.subheader("Notes on correlation analysis")
            correlcontent = st_quill(placeholder="Write your notes here", value = st.session_state.correlnotes, key="correlquill")

            st.session_state.correlnotes = correlcontent

            st.write("Correlation analysis took ", time.time() - start_time, "seconds to run")

        elif value == "Hypothesis testing":

            st.subheader("Hypothesis testing on characteristics with optional disaggregations by groups")
            
            with st.spinner('Please wait while we conduct hypothesis testing'):

                my_bar = st.progress(0)

                for percent_complete in range(100):
                    time.sleep(0.1)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                st.write("**Background to the model**")

                col1, col2, col3 = st.columns(3)
                with col1:
                    pass
                with col2:
                    st.image("ttest.png")
                with col3:
                    pass

                with st.expander("Assumptions of the one-sample t-test"):
                    st.success("One-sample t-test")
                    st.write("Student's t-test or t-test is a parametric inferential statistical method used for comparing the means between two different groups (two-sample t-test) or with the specific value (one-sample t-test).")
                    st.write("This test works well for hypothesis testing, particularly amongst data with low sample sizes (in contrast to the z-test).")
                    st.write("One Sample t-test (single sample t-test) is used to compare the sample mean (a random sample from a population) with the specific value (hypothesized or known mean of the population).")
                    st.info("Dependent variable should have an approximately normal distribution (Shapiro-Wilks Test).")
                    st.info("Observations are independent of each other.")
                    st.info("Note: One sample t-test is relatively robust to the assumption of normality when the sample size is large (n ≥ 30).")
                    st.warning("Null hypothesis: Sample mean is equal to the hypothesized or known population mean.")
                    st.warning("Alternative hypothesis: Sample mean is not equal to the hypothesized or known population mean (two-tailed or two-sided).")
                    st.warning("Alternative hypothesis: Sample mean is either greater or lesser to the hypothesized or known population mean (one-tailed or one-sided).")

                with st.expander("Assumptions of the two-sample t-test"):
                    st.success("One-sample t-test")
                    st.write("The two-sample (unpaired or independent) t-test compares the means of two independent groups, determining whether they are equal or significantly different. In two sample t-test, usually, we compute the sample means from two groups and derives the conclusion for the population’s means (unknown means) from which two groups are drawn.")
                    st.info("Observations in two groups have an approximately normal distribution (Shapiro-Wilks Test).")
                    st.info("Homogeneity of variances (variances are equal between treatment groups) (Levene or Bartlett Test).")
                    st.info("The two groups are sampled independently from each other from the same population.")
                    st.info("Note: Two sample t-test is relatively robust to the assumption of normality and homogeneity of variances when sample size is large (n ≥ 30) and there are equal number of samples (n1 = n2) in both groups.")
                    st.warning("Null hypothesis: Two group means are equal.")
                    st.warning("Alternative hypothesis: Two group means are different (two-tailed or two-sided).")
                    st.warning("Alternative hypothesis: Mean of one group either greater or lesser than another group (one-tailed or one-sided).")

                with st.expander("Sample size considerations for t-test"):
                    st.write("The t-test can be applied for the extremely small sample size (n = 2 to 5) provided the effect size is large and data follows the t-test assumptions. Remember, a larger sample size is preferred over small sample sizes.")
                    st.write("t-test is relatively robust to the assumption of normality and homogeneity of variances when the sample size is large (n ≥ 30).")

                # brands = pd.Series(df2['BRAND'].drop_duplicates()).sort_values().tolist() # Faster than .unique()
                # brandsmod = ["None"] + brands

                # Upload dataset
                st.title("Hypothesis Testing Tool")

                # Select the categorical variable of interest
                st.markdown("**Choose a categorical variable for grouping:**")
                categorical_vars = df.select_dtypes(include=['object', 'category']).columns.tolist()
                if not categorical_vars:
                    st.error("No categorical variables found in the dataset.")
                else:
                    group_var = st.selectbox("Select categorical variable:", categorical_vars, key="group_var")

                    # Get unique categories in the selected variable
                    categories = df[group_var].dropna().unique().tolist()
                    categories.insert(0, "None")

                    cat1 = st.selectbox("Select first category:", categories, key="cat1")
                    cat2 = st.selectbox("Select second category:", categories[1:], key="cat2")

                    # Select numerical variable for analysis
                    st.markdown("**Choose a numerical variable for analysis:**")
                    numerical_vars = df.select_dtypes(include=[np.number]).columns.tolist()
                    numerical_vars.insert(0, "NONE")
                    analysis_var = st.selectbox("Select numerical variable:", numerical_vars, key="analysis_var")

                    if analysis_var != "NONE":
                        if cat1 != "None":
                            # Filter data for the selected categories
                            dfa = df[df[group_var] == cat1]

                            if cat2 == "None":
                                # One-sample t-test
                                st.write("**Performing one-sample t-test:**")
                                sample_data = dfa[analysis_var].dropna()
                                population_mean = df[analysis_var].mean()
                                
                                if not sample_data.empty:
                                    t_stat, p_val = stats.ttest_1samp(sample_data, popmean=population_mean)
                                    st.write(f"T-statistic: {t_stat}, P-value: {p_val}")

                                    if p_val <= 0.1:
                                        st.success(f"P-value ({p_val}) is less than 0.1. Reject the null hypothesis.")
                                    else:
                                        st.error(f"P-value ({p_val}) is greater than 0.1. Fail to reject the null hypothesis.")
                                else:
                                    st.error("No data available for the selected category.")

                            else:
                                # Two-sample t-test
                                st.write("**Performing two-sample t-test:**")
                                dfb = df[df[group_var] == cat2]

                                sample_data_a = dfa[analysis_var].dropna()
                                sample_data_b = dfb[analysis_var].dropna()

                                if not sample_data_a.empty and not sample_data_b.empty:
                                    t_stat, p_val = stats.ttest_ind(sample_data_a, sample_data_b, equal_var=False)
                                    st.write(f"T-statistic: {t_stat}, P-value: {p_val}")

                                    if p_val <= 0.1:
                                        st.success(f"P-value ({p_val}) is less than 0.1. Reject the null hypothesis.")
                                    else:
                                        st.error(f"P-value ({p_val}) is greater than 0.1. Fail to reject the null hypothesis.")

                                    # Plot histogram
                                    st.write("**Histogram of selected variable by categories:**")
                                    histogram = df[df[group_var].isin([cat1, cat2])]
                                    histogram[group_var] = histogram[group_var].astype(str)

                                    plt.figure()
                                    for category in [cat1, cat2]:
                                        subset = histogram[histogram[group_var] == category]
                                        plt.hist(subset[analysis_var], alpha=0.5, label=str(category))
                                    plt.legend(title=group_var)
                                    plt.title(f"Histogram of {analysis_var} by {group_var}")
                                    plt.xlabel(analysis_var)
                                    plt.ylabel("Frequency")

                                    buf = BytesIO()
                                    plt.savefig(buf, format="png")
                                    buf.seek(0)
                                    st.image(buf)
                                else:
                                    st.error("No data available for one or both selected categories.")

                        else:
                            st.error("Please select at least one category.")

                    else:
                        st.error("Please select a numerical variable for analysis.")

            # Spawn a new Quill editor
            st.subheader("Notes on hypothesis analysis")
            hyptestingcontent = st_quill(placeholder="Write your notes here", key="hyptestingquill")

            st.session_state.hyptestingnotes = hyptestingcontent

            st.write("Hypothesis testing took ", time.time() - start_time, "seconds to run")

        elif value == "Dimension reduction":

            st.subheader("Dimension reduction")

            with st.spinner('Please wait while we set up the dimension reduction'):

                my_bar = st.progress(0)

                time.sleep(5)

                for percent_complete in range(100):
                    time.sleep(0.1)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                st.subheader("Dimension reduction using random forest analysis")

                # Handle missing or non-numeric values in the dataset
                st.write("**Cleaning the dataset:**")
                df.replace("?", np.nan, inplace=True)
                df.dropna(inplace=True)
                for col in df.select_dtypes(include=['object']).columns:
                    df[col], _ = pd.factorize(df[col])
                st.write("Cleaned Dataset Preview:")
                st.write(df.head())

                # Select the variable for dimension reduction
                st.markdown("**Choose a target variable for dimension reduction:**")
                dimredvars = ["None"] + list(df.columns)
                dimredvar = st.selectbox("Please choose variable of interest:", dimredvars, key="dimredy")

                if dimredvar != "None":
                    # Preparing target variable and features
                    y = df[dimredvar].values
                    X = df.drop(columns=[dimredvar])

                    # Train RandomForestClassifier
                    rf = RandomForestClassifier(n_estimators=100, max_depth=3,
                                                bootstrap=True, n_jobs=-1,
                                                random_state=0)
                    rf.fit(X, y)

                    # Feature importance
                    feature_imp = pd.Series(rf.feature_importances_, 
                                            index=X.columns).sort_values(ascending=False)

                    st.write('Feature importances:')
                    st.write(feature_imp)

                    # Display feature importance as a barplot
                    col1, col2, col3 = st.columns([0.5, 3, 0.5])
                    with col2:
                        sns.set(rc={'figure.figsize': (10, 8)})
                        plt.tick_params(axis='both', which='major', labelsize=5)
                        sns.barplot(x=feature_imp, y=feature_imp.index)
                        plt.xlabel('Feature Importance Score', fontsize=14)
                        plt.title(f"Visualizing Important Features for {dimredvar}", fontsize=18, pad=15)

                        buf = BytesIO()
                        plt.savefig(buf, format="png")
                        buf.seek(0)
                        st.image(buf)

                    # Feature selection based on importance threshold
                    st.markdown("**Feature Selection**")
                    chosenthreshold = st.number_input(
                        "Please select a threshold level of confidence to retain features (e.g., 0.05):", 
                        value=0.05, min_value=0.0, max_value=1.0, step=0.01
                    )

                    selector = SelectFromModel(rf, threshold=chosenthreshold, prefit=True)
                    features_important = selector.transform(X)

                    st.write('Sample of data with initial features:')
                    st.write(X.head())
                    st.info(f"Chosen level of significance to retain features: {chosenthreshold}")

                    reduceddf = pd.DataFrame(features_important, columns=X.columns[selector.get_support()])
                    st.session_state.reduceddf = reduceddf

                    st.write('Sample of data with selected features:')
                    st.write(reduceddf.head())

                    # Retrain RandomForest with reduced features
                    rf.fit(features_important, y)

                    numfeatures = len(reduceddf.columns)
                    featurestokeep = feature_imp.index[:numfeatures]

                    st.write('Retained features after thresholding:')
                    st.write(featurestokeep.tolist())

                    importantdf = df[featurestokeep.tolist() + [dimredvar]]
                    st.session_state.importantdf = importantdf

                    st.write('Sample of data with only important features retained per the chosen threshold:')
                    st.write(importantdf.head())

            # Spawn a new Quill editor
            st.subheader("Notes on dimension reduction analysis")
            dimredcontent = st_quill(placeholder="Write your notes here", key="dimredquill")

            st.session_state.dimrednotes = dimredcontent

            st.write("Dimension reduction took ", time.time() - start_time, "seconds to run")

        elif value == "ANOVA analysis":

            st.subheader("ANOVA analysis")

            st.warning("This model takes a long time to run. Please be patient.")

            with st.spinner('Please wait while we run the ANOVA analysis'):

                my_bar = st.progress(0)

                time.sleep(10)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                with st.expander("Explanation of ANOVA analysis"):
                    st.write("Analysis of Variance (ANOVA) is used to test how different sample groups within a dataset compare to each other. In our case, it will be used to test either (a) how one (or a group) of categories compare to the overall dataset, or (b) how two different groups of categories compare to each other, with respect to a chosen variable.")
                    st.write("Groups mean differences are inferred by analyzing variances. ANOVA uses variance-based F test to check the group mean equality.")
                    st.info("The null hypothesis that group means are equal is tested with an F-test for all groups, followed by post-hoc tests to see individual group differences.")

                # Replace '?' with NaN and drop missing values
                df.replace('?', np.nan, inplace=True)
                df.dropna(inplace=True)
                st.write("Cleaned Dataset preview:")
                st.dataframe(df.head())

                # Allow user to select categorical variables and interaction variables
                categorical_vars = df.select_dtypes(include=['object', 'category']).columns.tolist()
                if not categorical_vars:
                    st.error("No categorical variables found in the dataset.")
                else:
                    anova_var = st.selectbox("Select a categorical variable for ANOVA analysis:", categorical_vars)
                    interaction_var = st.selectbox("Select an interaction variable (optional):", ["None"] + categorical_vars)

                    # Allow user to select comparison groups
                    unique_groups = df[anova_var].unique().tolist()
                    unique_groups_mod = ["All"] + unique_groups

                    st.markdown("**Please choose categories of interest.**")
                    group_selection = st.multiselect("Select categories of interest:", unique_groups_mod, key="groups")

                    st.markdown("**Please choose second categories of interest for comparison.**")
                    group_selection2 = st.multiselect("Select second (comparison) categories:", unique_groups, key="groups2")

                    # Allow user to select the dependent variable
                    numerical_vars = df.select_dtypes(include=['float64', 'int64']).columns.tolist()
                    dependent_var = st.selectbox("Select the dependent variable:", numerical_vars)

                    if group_selection and dependent_var:
                        filtered_data = df.copy()
                        
                        if "All" not in group_selection:
                            filtered_data = filtered_data[filtered_data[anova_var].isin(group_selection)]

                        if group_selection2:
                            filtered_data2 = df[df[anova_var].isin(group_selection2)]
                            combined_data = pd.concat([filtered_data, filtered_data2])
                        else:
                            combined_data = filtered_data

                        st.write("Filtered Dataset Preview:")
                        st.dataframe(combined_data.head())

                        try:
                            # Build ANOVA model
                            if interaction_var == "None":
                                anova_model = f'Q("{dependent_var}") ~ C(Q("{anova_var}"))'
                                st.info("Running ANOVA without interaction variable.")
                            else:
                                anova_model = f'Q("{dependent_var}") ~ C(Q("{anova_var}")) + C(Q("{interaction_var}")) + C(Q("{anova_var}")):C(Q("{interaction_var}"))'
                                st.info("Running ANOVA with interaction variable.")

                            # Print the formula for debugging
                            st.write(f"ANOVA Formula: {anova_model}")

                            # ANOVA analysis
                            res = stat()
                            res.anova_stat(df=combined_data, res_var=dependent_var, anova_model=anova_model)
                            st.write(res.anova_summary)

                            st.success("If the p-value (PR(>F)) is smaller than 0.05, we can reject the null hypothesis and conclude significant differences between groups.")

                            # Tukey HSD test for post-hoc comparisons
                            res.tukey_hsd(df=combined_data, res_var=dependent_var, xfac_var=anova_var, anova_model=anova_model)
                            st.write(res.tukey_summary)

                            # Visualizations
                            st.info("Visualizing results.")
                            if interaction_var == "None":
                                sns.boxplot(x=anova_var, y=dependent_var, data=combined_data, palette="Set3")
                                sns.swarmplot(x=anova_var, y=dependent_var, data=combined_data, color="black", size=3)
                            else:
                                sns.boxplot(x=anova_var, y=dependent_var, hue=interaction_var, data=combined_data, palette="Set3")
                            st.pyplot()

                            # Interaction plot
                            if interaction_var != "None":
                                st.info("Visualizing interaction effect.")
                                combined_data = combined_data.reset_index(drop=True)

                                # Generate colors dynamically based on interaction variable levels
                                interaction_levels = combined_data[interaction_var].nunique()
                                color_palette = plt.cm.tab10(np.linspace(0, 1, interaction_levels))  # Generate enough colors

                                # Generate interaction plot
                                try:
                                    fig = interaction_plot(
                                        combined_data[anova_var],
                                        combined_data[interaction_var],
                                        combined_data[dependent_var],
                                        colors=color_palette
                                    )
                                    st.pyplot(fig)
                                except Exception as e:
                                    st.error(f"Error generating interaction plot: {e}")

                            # Residual analysis
                            st.info("Checking assumptions: Normality and Homogeneity of Variances")

                            # Shapiro-Wilk test for normality
                            w, pvalue = stats.shapiro(res.anova_model_out.resid)
                            st.write("Shapiro-Wilk Test (Normality): W = {:.3f}, p-value = {:.3f}".format(w, pvalue))
                            if pvalue < 0.05:
                                st.error("The residuals are not normally distributed. Interpret results cautiously.")
                            else:
                                st.success("The residuals are normally distributed.")

                            # Bartlett's test for homogeneity of variances
                            bartlett_w, bartlett_pvalue = stats.bartlett(*[group[dependent_var].values for name, group in combined_data.groupby(anova_var)])
                            st.write("Bartlett's Test (Equal Variances): W = {:.3f}, p-value = {:.3f}".format(bartlett_w, bartlett_pvalue))
                            if bartlett_pvalue < 0.05:
                                st.error("Variances are not equal across groups. Consider data transformations.")
                            else:
                                st.success("Variances are equal across groups.")

                            # Levene's test
                            res.levene(df=combined_data, res_var=dependent_var, xfac_var=anova_var)
                            st.write(res.levene_summary)

                            # QQ plot
                            sm.qqplot(res.anova_std_residuals, line='45')
                            plt.xlabel("Theoretical Quantiles")
                            plt.ylabel("Standardized Residuals")
                            st.pyplot()

                            # Histogram of residuals
                            plt.hist(res.anova_model_out.resid, bins='auto', histtype='bar', ec='k')
                            plt.xlabel("Residuals")
                            plt.ylabel("Frequency")
                            st.pyplot()

                            st.success("ANOVA analysis completed successfully.")

                            # Add explanations of results
                            st.subheader("Interpreting the Results")

                            # Dynamic explanation based on ANOVA results
                            if res.anova_summary['PR(>F)'][0] < 0.05:
                                st.write(f"The ANOVA results show that the effect of {anova_var} on {dependent_var} is statistically significant (p < 0.05). This means there is evidence to suggest that at least one group mean differs from the others.")
                            else:
                                st.write(f"The ANOVA results show that the effect of {anova_var} on {dependent_var} is not statistically significant (p >= 0.05). This means we do not have enough evidence to conclude that the group means differ.")

                            if interaction_var != "None":
                                if res.anova_summary['PR(>F)'][2] < 0.05:
                                    st.write(f"The interaction between {anova_var} and {interaction_var} is statistically significant (p < 0.05). This suggests that the combined effect of these variables influences {dependent_var}.")
                                else:
                                    st.write(f"The interaction between {anova_var} and {interaction_var} is not statistically significant (p >= 0.05). This suggests that their combined effect does not significantly influence {dependent_var}.")

                            # Dynamic explanation based on post-hoc Tukey results
                            try:
                                significant_pairs = res.tukey_summary[res.tukey_summary['p-adj'] < 0.05]
                            except KeyError:
                                if 'p_value' in res.tukey_summary.columns:
                                    significant_pairs = res.tukey_summary[res.tukey_summary['p_value'] < 0.05]
                                elif 'adj_pval' in res.tukey_summary.columns:
                                    significant_pairs = res.tukey_summary[res.tukey_summary['adj_pval'] < 0.05]
                                else:
                                    significant_pairs = pd.DataFrame()  # Default to empty if no valid p-value column exists

                            if not significant_pairs.empty:
                                st.write("The following pairwise comparisons show significant differences:")
                                st.dataframe(significant_pairs)
                            else:
                                st.write("No significant pairwise differences in the Tukey HSD test were detected, or p-value information is unavailable.")

                        except Exception as e:
                            st.error(f"Error during ANOVA analysis: {e}")

            # Spawn a new Quill editor
            st.subheader("Notes on ANOVA analysis")
            anovacontent = st_quill(placeholder="Write your notes here", key="anovaquill")

            st.session_state.anovanotes = anovacontent

            st.write("ANOVA analysis took {:.2f} seconds to run".format(time.time() - start_time))
                            
        elif value == "Regression analysis":

            st.subheader("Linear regression analysis")

            with st.spinner('Please wait while we conduct the linear regression analysis'):

                my_bar = st.progress(0)

                time.sleep(5)

                for percent_complete in range(100):
                    time.sleep(0.2)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                ### From Jupyter - 2. Linear regression

                # Choose predicted variable - this will become dynamic in the app
                regvars = list(df.columns)
                # regvars.insert(0, "NONE")
                # y = df['BRAND'] # old predicted variable
                chosenvar = st.selectbox("Please select variable for analysis", ["Choose dependent variable"] + regvars)

                if chosenvar != "Choose dependent variable":

                    y = df[chosenvar]
                    df3 = y.copy()

                    # Define predictor variables
                    df2 = df.drop([chosenvar], axis=1)
                    indepvars = st.multiselect("Please select independent (input) variables for analysis", regvars)
                    # df3 = df2.copy()
                    # df3 = y.copy()

                    if indepvars != []:

                        for indepvar in indepvars:
                            dfi = df[indepvar]
                            df3 = pd.concat((df3, dfi), axis=1)
                        x = df3.drop([chosenvar], axis=1)
                        # print(df3)
                        # x = df3.iloc[:, :]

                    if chosenvar != "NONE" and indepvars != []: 

                        # x, y = np.array(x), np.array(y)

                        x = sm.add_constant(x)

                        model = sm.OLS(y, x)

                        results = model.fit()

                        st.write(results.summary())

                        # Use the DataFrame df3 to get column names
                        coef_df = pd.DataFrame({
                            "Variable": ["const"] + list(df3.drop([chosenvar], axis=1).columns),
                            "Coefficient": results.params
                        }).sort_values(by="Coefficient", ascending=False)

                        # Plot the regression coefficients
                        fig, ax = plt.subplots(figsize=(10, 6))
                        sns.barplot(x="Coefficient", y="Variable", data=coef_df, ax=ax)
                        plt.title("Regression Coefficients")
                        buf = BytesIO()
                        plt.savefig(buf, format="png")
                        st.image(buf)

                        st.write("")

                        st.write('Predicted response:', results.fittedvalues, sep='\n') # Or print('predicted response:', results.predict(x), sep='\n')

                        ### End

                        st.write("")

                        with st.expander("Interpretation guide"):
                            st.write("Regression searches for relationships (mathematical dependencies) among variables. This differs from correlation because regression analysis seeks to unpack **causal** relationships, where a change in one variable causes another variable to change. We find a function that maps some features or variables to others sufficiently well.\n\n"
                            "The dependent features are called the dependent variables, outputs, or responses. The independent features are called the independent variables, inputs, or predictors.\n\n"
                            "Regression problems usually have one continuous and unbounded dependent variable. The inputs, however, can be continuous, discrete, or even categorical data such as gender, nationality, brand, and so on. It is a common practice to denote the outputs with y and inputs with x.\n\n"
                            "Regression is also useful when you want to forecast a response using a new set of predictors. The regression coefficients can be applied as a multiplier factor to understand how changes in an predictor feature would cause the predicted variable to change.")
                            st.info("Regression coefficients indicate the predicted change in the output variable caused by a change in the input variable.")
                            st.info("The estimated or predicted response for each observation should be as close as possible to the corresponding actual response in the underlying data (in real life). The differences for all observations from the actual responses are called the residuals. Regression is about determining the best predicted weights, that is the weights corresponding to the smallest residuals. To get the best weights, you usually minimize the sum of squared residuals (SSR) for all observations. This approach is called the method of ordinary least squares.")
                            st.info("The variation of actual responses occurs partly due to the dependence on the predictors. However, there is also an additional, systematic, inherent variance of the output. The coefficient of determination, denoted as R-squared, tells you which amount of variation in 𝑦 can be explained by the dependence on 𝐱 using the particular regression model. Larger 𝑅² indicates a better fit and means that the model can better explain the variation of the output with different inputs.")
                            st.error("You should, however, be aware of two problems that might follow the choice of the degree: underfitting and overfitting.\n"
                            "Underfitting occurs when a model can't accurately capture the dependencies among data, usually as a consequence of its own simplicity. It often yields a low 𝑅² with known data and bad generalization capabilities when applied with new data. Overfitting happens when a model learns both dependencies among data and random fluctuations. In other words, a model learns the existing data too well. Complex models, which have many features or terms, are often prone to overfitting. When applied to known data, such models usually yield high 𝑅². However, they often don't generalize well and have significantly lower 𝑅² when used with new data.")
                            st.success("The coefficients in our model indicate the extent to which we can expect the predicted variable chosen to change for a 1 unit change in each respective input variable. The p-values indicate the level of statistical significance (lower means more statistically significant). R-squared is a measure of how well the model explains variance in the data.")

            # Spawn a new Quill editor
            st.subheader("Notes on linear regression analysis")
            regcontent = st_quill(placeholder="Write your notes here", key="regquill")

            st.session_state.regnotes = regcontent

            st.write("Linear regression took ", time.time() - start_time, "seconds to run")

        elif value == "Clustering analysis":

            st.subheader("Cluster analysis")

            st.info("We run cluster analysis for any variable, with optional grouping by brand, for a more granular analysis, if desired.")

            with st.spinner('Please wait while we conduct the cluster analysis using the K-means algorithm'):

                my_bar = st.progress(0)

                time.sleep(10)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                # Handle missing or non-numeric values in the dataset
                st.write("**Cleaning the dataset:**")
                df.replace("?", np.nan, inplace=True)
                df.dropna(inplace=True)
                for col in df.select_dtypes(include=['object']).columns:
                    df[col], _ = pd.factorize(df[col])
                st.write("Cleaned Dataset Preview:")
                st.write(df.head())

                # Assumptions and operation of the model
                with st.expander("Assumptions and operation of the model"):
                    st.write("The technique of cluster analysis aims to group similar observations in a dataset, such that observations in the same group are as similar to each other as possible, and similarly, observations in different groups are as different to each other as possible.")
                    st.warning("This allows us to get a sense of how many different groups with different characteristics and observable behaviour we have in our data, with respect to a chosen variable.")
                    st.info("Cluster analysis does this by minimizing the distance between observations (or the error involved in grouping them), by considering the mean of each potential group. Each group, by construct, has specific characteristics as it relates to the chosen variable.")
                    st.write("Cluster analysis in this case assumes each explanatory variable has the same within-group variance, with spherical variance, and that clusters are roughly similarly sized. Our data is sourced in a way that produces usable variance, and cluster size as reported can be used to further refine results by choosing different cluster sizes as a tuned hyperparameter.")

                # Select the variable for cluster analysis
                st.markdown("**Choose a variable for cluster analysis:**")
                clustervars = list(df.columns)
                chosenvar = st.selectbox("Please select variable for clustering:", ["None"] + clustervars)

                # Optional grouping by a categorical variable
                st.markdown("**Optionally group data by a categorical variable:**")
                st.info("Tip: it is instructive to compare analysis between categories of your chosen variable for different interactions from other variables.")
                cat_vars = df.select_dtypes(include=['int64', 'category', 'object']).columns.tolist()
                group_var = st.selectbox("Select a grouping variable (optional):", ["None"] + cat_vars)

                if chosenvar != "None":
                    X2 = df[[chosenvar]].copy()

                    if group_var != "None":
                        st.markdown("**Filter data by selected groups:**")
                        groups = df[group_var].unique()
                        selected_groups = st.multiselect("Select groups to include in the analysis:", groups, default=groups)
                        X2 = df[df[group_var].isin(selected_groups)][[chosenvar]]

                    # Find optimal number of clusters
                    st.markdown("**Optimal Number of Clusters:**")

                    # Elbow Method
                    distortions = []
                    for i in range(1, 16):
                        km = KMeans(n_clusters=i, init='k-means++', n_init=10, max_iter=300, tol=1e-4, random_state=0)
                        km.fit(X2)
                        distortions.append(sum(np.min(cdist(X2, km.cluster_centers_, 'euclidean'), axis=1)) / X2.shape[0])

                    # Plot distortions
                    fig, ax = plt.subplots(figsize=(10, 6))
                    plt.plot(range(1, 16), distortions, marker='o')
                    plt.xlabel('Number of clusters')
                    plt.ylabel('Distortion')
                    buf = BytesIO()
                    fig.savefig(buf, format="png")
                    st.image(buf)

                    st.info("Use the elbow method graph to determine the optimal number of clusters where the curve starts to level off. You can use the elbow method to refine the analysis, by informing the number of clusters to use below. You should choose the number of clusters based on where the graph sensibly starts to taper/level off, meaning that additional clusters add little additional explanatory power to the model.")

                    # Silhouette Method
                    sil = []
                    kmax = 10
                    nclusters = 0

                    for k in range(2, kmax + 1):
                        kmeans = KMeans(n_clusters=k, random_state=0).fit(X2)
                        labels = kmeans.labels_
                        silscore = silhouette_score(X2, labels, metric='euclidean')
                        sil.append(silscore)
                        if silscore == max(sil):
                            nclusters = k

                    # Plot silhouette scores
                    fig, ax = plt.subplots(figsize=(10, 6))
                    plt.plot(range(2, kmax + 1), sil, marker='o')
                    plt.xlabel('Number of clusters')
                    plt.ylabel('Silhouette Score')
                    buf = BytesIO()
                    fig.savefig(buf, format="png")
                    st.image(buf)

                    st.info(f"The optimal number of clusters based on the silhouette score is {nclusters}.")

                    # Button to run clustering with maximum silhouette score
                    silrun = st.button("Run cluster analysis with maximum silhouette score")

                    if silrun:
                        kmnew = KMeans(n_clusters=nclusters, init='k-means++', n_init=10, max_iter=300, tol=1e-4, random_state=0)
                        kmnew.fit(X2)
                        X2['CLUSTERS'] = kmnew.labels_

                        for i in range(nclusters):
                            countclusters = X2[X2['CLUSTERS'] == i].shape[0]
                            st.write(f"Number of Rows in cluster {i}: {countclusters}")

                        # Plot the following variables and their clusters
                        fig = sns.pairplot(X2, vars=[chosenvar], hue="CLUSTERS", palette=sns.color_palette("hls", nclusters), height=5)
                        buf = BytesIO()
                        fig.savefig(buf, format="png")
                        st.image(buf)

                        st.success(f"Cluster analysis completed with {nclusters} clusters. Each cluster represents a distinct group of data for the variable '{chosenvar}'. Further analysis can help assign meaning to these clusters.")

                    st.write("Or choose the number of clusters based on the elbow method, enter this in the slider below, and click the 'Conduct analysis' button")
                    numclus = st.slider("Please choose number of clusters", min_value=1, max_value=10)
                    ownclus = st.button("Conduct analysis")
                    if ownclus:
                        kmnew = KMeans(n_clusters=numclus, init='k-means++', n_init=10, max_iter=300, tol=1e-4, random_state=0)
                        kmnew.fit(X2)
                        X2['CLUSTERS'] = kmnew.labels_

                        for i in range(numclus):
                            countclusters = X2[X2['CLUSTERS'] == i].shape[0]
                            st.write(f"Number of Rows in cluster {i}: {countclusters}")

                        # Plot the following variables and their clusters
                        fig = sns.pairplot(X2, vars=[chosenvar], hue="CLUSTERS", palette=sns.color_palette("hls", numclus), height=5)
                        buf = BytesIO()
                        fig.savefig(buf, format="png")
                        st.image(buf)

                        st.success("You chose " + str(numclus) + " clusters in the data for " + chosenvar + " for the brands selected. Each group has specific characteristics and behaviour that be intuited (you can assign meanings to each group based on hypotheses), and further investigation using data exploration, correlation analysis, and linear/logistic regression analysis can provide further insight for " + chosenvar + " clusters in these brands, as compared to other brands or the whole dataset.")
                        
                # Optional - log transformations

                ### End

            st.write("")

            st.subheader("Notes on cluster analysis")

            # Spawn a new Quill editor
            st.subheader("Notes on cluster analysis")
            clustercontent = st_quill(placeholder="Write your notes here", key="clusterquill")

            st.session_state.clusternotes = clustercontent

            st.write("Running the cluster analysis using K-means took ", time.time() - start_time, "seconds to run")

        elif value == "Conjoint analysis":

            st.subheader("Conjoint analysis")

            st.warning("This model takes some time to run. Please be patient.")

            with st.spinner('Please wait while we conduct the conjoint analysis'):

                my_bar = st.progress(0)

                time.sleep(10)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                ### From Jupyter - conjoint analysis

                st.info("Conjoint analysis is traditionally a method for determining what features of a product are most important to consumers. We extend the usage here to enable us to determine what factors are most important in determining an individual's outcomes related to a chosen variable.")

                max_elements = 1000000  # Set the maximum allowed elements for Styler
                pd.set_option("styler.render.max_elements", max_elements)
                st.dataframe(df.style.highlight_null(null_color='red'))

                # Select categorical variable for grouping
                cat_vars = df.select_dtypes(include=['object', 'category']).columns.tolist()
                if not cat_vars:
                    st.error("No categorical variables available for grouping.")
                else:
                    group_var = st.selectbox("Select a categorical variable for grouping:", cat_vars)

                    if group_var:
                        # Get unique categories dynamically
                        categories = pd.Series(df[group_var].drop_duplicates()).sort_values().tolist()

                        sel_category = st.selectbox("Please select a category to analyze", categories)

                        # Subset the data based on the selected category
                        category_code = df.loc[df[group_var] == sel_category, group_var].iloc[0]

                        y = df[group_var].apply(lambda x: 1 if x == category_code else 0).head(10000)  # Subsampling for speed and efficiency
                        x = df[[col for col in df.columns if col != group_var]].head(10000)  # Subsampling for speed and efficiency

                        # One-hot encode all categorical variables except the grouping variable
                        xdum = pd.get_dummies(x, columns=[c for c in x.select_dtypes(include=['object', 'category']).columns if c != group_var])
                        st.write("Transformed Data (First 5 Rows):")
                        st.dataframe(xdum.head())

                        plt.style.use('bmh')

                        # Perform regression
                        try:
                            res = sm.OLS(y.astype(float), xdum.astype(float)).fit()
                            st.write(res.summary().as_text()[:1000] + "...")
                            if st.button("View Full Summary"):
                                st.text(res.summary())
                        except Exception as e:
                            st.error(f"An error occurred during regression: {e}")

                        st.subheader("Factor importances:")

                        df_res = pd.DataFrame({
                        'param_name': res.params.keys()
                        , 'param_w': res.params.values
                        , 'pval': res.pvalues
                        })
                        # adding field for absolute of parameters
                        df_res['abs_param_w'] = np.abs(df_res['param_w'])
                        # marking field is significant under 95% confidence interval
                        df_res['is_sig_95'] = (df_res['pval'] < 0.05)
                        # constructing color naming for each param
                        df_res['c'] = ['blue' if x else 'red' for x in df_res['is_sig_95']]

                        # make it sorted by abs of parameter value
                        df_res = df_res.sort_values(by='abs_param_w', ascending=True)

                        fig, ax = plt.subplots(figsize=(14, 8))
                        plt.title('Factor Importances')
                        pwu = df_res['param_w']
                        xbar = np.arange(len(pwu))
                        plt.barh(xbar, pwu, color=df_res['c'])
                        plt.yticks(xbar, labels=[label[:15] + '...' if len(label) > 15 else label for label in df_res['param_name']])

                        buf = BytesIO()
                        fig.savefig(buf, format="png")
                        st.image(buf)

                        st.subheader("Absolute and relative/normalized importances:")

                        # Compute feature importance
                        range_per_feature = dict()
                        for key, coeff in res.params.items():
                            sk = key.split('_')
                            feature = sk[0]
                            if len(sk) == 1:
                                feature = key
                            if feature not in range_per_feature:
                                range_per_feature[feature] = list()

                            range_per_feature[feature].append(coeff)

                        importance_per_feature = {
                            k: max(v) - min(v) for k, v in range_per_feature.items()
                        }

                        total_feature_importance = sum(importance_per_feature.values())
                        relative_importance_per_feature = {
                            k: 100 * round(v/total_feature_importance, 3) for k, v in importance_per_feature.items()
                        }

                        alt_data = pd.DataFrame(
                            list(importance_per_feature.items()), 
                            columns=['attr', 'importance']
                        ).sort_values(by='importance', ascending=False)

                        fig, ax = plt.subplots(figsize=(12, 8))
                        xbar = np.arange(len(alt_data['attr']))
                        plt.title('Importance')
                        plt.barh(xbar, alt_data['importance'])
                        for i, v in enumerate(alt_data['importance']):
                            ax.text(v, i + .25, '{:.2f}'.format(v))
                        plt.ylabel('Attributes')
                        plt.xlabel('% Importance')
                        plt.yticks(xbar, alt_data['attr'])

                        buf = BytesIO()
                        fig.savefig(buf, format="png")
                        st.image(buf)

                        alt_data = pd.DataFrame(
                            list(relative_importance_per_feature.items()), 
                            columns=['attr', 'relative_importance (pct)']
                        ).sort_values(by='relative_importance (pct)', ascending=False)

                        fig, ax = plt.subplots(figsize=(12, 8))
                        xbar = np.arange(len(alt_data['attr']))
                        plt.title('Relative Importance / Normalized Importance')
                        plt.barh(xbar, alt_data['relative_importance (pct)'])
                        for i, v in enumerate(alt_data['relative_importance (pct)']):
                            ax.text(v, i + .25, '{:.2f}%'.format(v))
                        plt.ylabel('Attributes')
                        plt.xlabel('% Relative Importance')
                        plt.yticks(xbar, alt_data['attr'])

                        buf = BytesIO()
                        fig.savefig(buf, format="png")
                        st.image(buf)

                        st.write("Top 3 Takeaways:")
                        top_takeaways = alt_data.head(3)
                        for i, row in top_takeaways.iterrows():
                            st.markdown(f"- **{row['attr']}**: {row['relative_importance (pct)']:.2f}%")

                        st.success("The results of the conjoint analysis show us which factors should be paid particular attention, in framing an individuals's choices relating to the chosen variable, " + group_var + ".")

            default_notes = """Top contributing factors:
                1. __
                2. __
                3. __

                Potential implications:
                - __"""

            # Spawn a new Quill editor
            st.subheader("Notes on conjoint analysis")
            conjointcontent = st_quill(placeholder=default_notes, key="conjointquill")

            st.session_state.conjointnotes = conjointcontent

            st.write("Running the conjoint analysis took ", time.time() - start_time, "seconds to run")

        elif value == "Neural networks":

            st.subheader("Neural Network Analysis")

            st.warning("This model can take a long time to run. Please be patient.")

            with st.spinner('Please wait while we conduct the neural networks analysis'):

                my_bar = st.progress(0)

                time.sleep(10)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                st.info("This neural networks model examines the effect of selected features on a target variable. It is designed for user-uploaded datasets and supports feature selection and model optimization.")

                with st.expander("What are variables for analysis and target variables?"):
                    st.write("**Variables for Analysis:** These are the features or predictors used by the model to make predictions. They can include numerical variables like age or income and categorical variables like gender or city.")
                    st.write("**Target Variable:** This is the outcome the model is trying to predict. It is typically a categorical variable like brand preference or customer segment.")

                # Handle missing values
                df.replace('?', np.nan, inplace=True)
                df.dropna(inplace=True)
                st.write("Cleaned Dataset preview:")
                st.dataframe(df.head())

                # Check data variability and distribution
                st.subheader("Data Distribution and Variability Check")
                for column in df.columns:
                    st.write(f"**{column}:**")
                    st.bar_chart(df[column].value_counts())
                    if df[column].nunique() == 1:
                        st.warning(f"The column `{column}` has only one unique value. This lack of variability might limit its predictive power.")
                    elif df[column].nunique() < 5:
                        st.info(f"The column `{column}` has low variability with only {df[column].nunique()} unique values.")

                # Allow user to select features and target variable
                numerical_vars = df.select_dtypes(include=['float64', 'int64']).columns.tolist()
                categorical_vars = df.select_dtypes(include=['object', 'category']).columns.tolist()

                if not numerical_vars or not categorical_vars:
                    st.error("Dataset must contain both numerical and categorical variables.")
                else:
                    st.markdown("**Select the variables to use in your analysis:**")
                    st.write("Choose features that are relevant predictors for your analysis, such as age, income, or city.")
                    feature_vars = st.multiselect("Select features for analysis:", numerical_vars + categorical_vars)

                    st.markdown("**Select the target variable:**")
                    st.write("Choose the variable you want to predict. For example, this could be a brand choice or customer category.")
                    target_var = st.selectbox("Select the target variable:", categorical_vars)

                    if feature_vars and target_var:
                        # Prepare data for analysis
                        X = pd.get_dummies(df[feature_vars], drop_first=True)
                        y = df[target_var]

                        # Encode target variable
                        le = LabelEncoder()
                        y = le.fit_transform(y)

                        # Ensure feature names are consistent
                        feature_names = X.columns.tolist()

                        # Split the data
                        X_train, X_test, y_train, y_test = train_test_split(X, y, random_state=0)

                        # Scale the data
                        scaler = StandardScaler()
                        X_train = scaler.fit_transform(X_train)
                        X_test = scaler.transform(X_test)

                        # Create and train the neural network model
                        reg = MLPClassifier(max_iter=2000, hidden_layer_sizes=(5, 5), random_state=1)
                        reg.fit(X_train, y_train)

                        # Evaluate initial model
                        y_pred = reg.predict(X_test)
                        accuracy = accuracy_score(y_test, y_pred)
                        st.write("Initial Accuracy: {:.3f}".format(accuracy))

                        # Optimize the model
                        validation_scores = {}
                        for hidden_layer_size in [(i, j) for i in range(3, 7) for j in range(3, 7)]:
                            reg = MLPClassifier(max_iter=2000, hidden_layer_sizes=hidden_layer_size, random_state=1)
                            score = cross_val_score(estimator=reg, X=X_train, y=y_train, cv=2)
                            validation_scores[hidden_layer_size] = score.mean()

                        # Find the best parameters
                        best_score = max(validation_scores.values())
                        optimal_hidden_layer_size = [k for k, v in validation_scores.items() if v == best_score][0]
                        st.write("Best Validation Score: {:.4f}".format(best_score))
                        st.write("Optimal Hidden Layer Sizes: {}".format(optimal_hidden_layer_size))

                        # Retrain with the best parameters
                        clf = MLPClassifier(max_iter=2000, hidden_layer_sizes=optimal_hidden_layer_size, random_state=1)
                        clf.fit(X_train, y_train)
                        y_pred_optimized = clf.predict(X_test)
                        optimized_accuracy = accuracy_score(y_test, y_pred_optimized)
                        st.write("Optimized Accuracy: {:.3f}".format(optimized_accuracy))

                        # Feature Impact Analysis
                        st.subheader("Feature Impact Analysis")
                        variable_impacts = []
                        for feature in feature_vars:
                            st.write(f"Analyzing the impact of {feature} on {target_var}...")
                            feature_values = sorted(df[feature].unique())
                            predictions = []

                            for value in feature_values:
                                X_design = pd.DataFrame(columns=feature_names)
                                for col in feature_names:
                                    X_design[col] = [X[col].median()] * len(feature_values)

                                if feature in X_design.columns:
                                    X_design[feature] = value

                                X_scaled = scaler.transform(X_design)
                                pred = clf.predict(X_scaled)
                                predictions.append(pred.mean())

                            variable_impacts.append({
                                "feature": feature,
                                "feature_values": feature_values,
                                "predictions": predictions
                            })

                            plt.plot(feature_values, predictions)
                            plt.xlabel(feature)
                            plt.ylabel(target_var)
                            plt.title(f"Effect of {feature} on {target_var}")
                            st.pyplot()

                        st.subheader("Interpreting the Results")

                        # Explain results dynamically
                        if optimized_accuracy > accuracy:
                            st.success(f"Model optimization improved accuracy from {accuracy:.3f} to {optimized_accuracy:.3f}.")
                        else:
                            st.warning(f"Model optimization did not improve accuracy. Initial accuracy: {accuracy:.3f}, Optimized accuracy: {optimized_accuracy:.3f}.")

                        st.write("Based on the feature impact analysis:")
                        for impact in variable_impacts:
                            feature = impact['feature']
                            feature_values = impact['feature_values']
                            predictions = impact['predictions']

                            if len(set(predictions)) > 1:
                                highest_effect = max(predictions)
                                lowest_effect = min(predictions)
                                highest_value = feature_values[predictions.index(highest_effect)]
                                lowest_value = feature_values[predictions.index(lowest_effect)]

                                impact_difference = highest_effect - lowest_effect

                                st.write(f"- **{feature}:** Changing `{feature}` from `{lowest_value}` to `{highest_value}` results in a `{impact_difference:.2f}` change in the predicted mean value of `{target_var}`. The highest impact occurs when `{feature}` is `{highest_value}`, with a predicted mean value of `{highest_effect:.2f}`. The lowest impact occurs when `{feature}` is `{lowest_value}`, with a predicted mean value of `{lowest_effect:.2f}`.")
                            else:
                                st.write(f"- **{feature}:** No significant variation in predictions for `{feature}`. This could be due to limited variability or insufficient relationship with `{target_var}`.")

                        st.write("Neural network analysis is complete.")

            # Spawn a new Quill editor
            st.subheader("Notes on neural networks analysis")
            neuralnetworkscontent = st_quill(placeholder="Write your notes here", key="neuralnetworksquill")

            st.session_state.neuralnetworksnotes = neuralnetworkscontent

            st.write("Neural network analysis took {:.2f} seconds to run.".format(time.time() - start_time))

        elif value == "Decision trees":

            st.subheader("Decision tree analysis")

            with st.spinner('Please wait while we conduct the decision tree analysis'):

                my_bar = st.progress(0)

                time.sleep(10)

                for percent_complete in range(100):
                    time.sleep(0.01)
                    my_bar.progress(percent_complete + 1)

                start_time = time.time()

                st.info("A decision tree is a decision support tool that uses a tree-like model of decisions and their possible consequences, including chance event outcomes, resource costs, and utility. We use a Decision Tree Regressor model, which considers causality in the model.")
                st.warning("In this way, a decision tree shows how a particular decision was made using selected explanatory variables.")

                # Replace '?' with NaN
                df.replace('?', np.nan, inplace=True)

                # Drop rows with missing values
                df.dropna(inplace=True)
                st.write("Cleaned Dataset preview:")
                st.dataframe(df.head())

                # Encode categorical variables
                categorical_cols = df.select_dtypes(include=['object', 'category']).columns
                label_encoders = {}
                for col in categorical_cols:
                    le = LabelEncoder()
                    df[col] = le.fit_transform(df[col].astype(str))
                    label_encoders[col] = le

                # Allow user to select the target variable
                target_var = st.selectbox("Select the target variable:", df.columns.tolist())

                # Allow user to select explanatory variables
                explanatory_vars = st.multiselect("Select 3 explanatory variables:", df.columns.tolist())

                if len(explanatory_vars) != 3:
                    st.error("Please select exactly 3 explanatory variables.")
                else:
                    # Prepare data for the decision tree
                    X = df[explanatory_vars]
                    y = df[target_var]

                    # Ensure no ambiguity in index
                    X = X.reset_index(drop=True)
                    y = y.reset_index(drop=True)

                    # Fit data to tree-based regression model
                    regressor = DecisionTreeRegressor(random_state=0)
                    regressor = regressor.fit(X, y)

                    # Visualizing the decision tree regression results
                    plt.figure(figsize=(6, 6), dpi=150)
                    plot_tree(regressor, max_depth=3, feature_names=X.columns, impurity=False, filled=True)
                    plt.show()

                    buf = BytesIO()
                    plt.savefig(buf, format="png")
                    st.image(buf)

                    # Scatter plot of first two explanatory variables
                    plt.figure(figsize=[6, 4], dpi=120)
                    plt.xlabel(explanatory_vars[0])
                    plt.ylabel(explanatory_vars[1])
                    plt.scatter(x=X[explanatory_vars[0]], y=X[explanatory_vars[1]], c=y, cmap='viridis', s=20, alpha=0.8)
                    plt.colorbar(label=target_var)
                    plt.show()

                    buf2 = BytesIO()
                    plt.savefig(buf2, format="png")
                    st.image(buf2)

                    # Generate synthetic data for visualization of decision boundaries
                    X_synthetic, y_synthetic = make_regression(n_samples=1000, n_features=2, n_informative=2, random_state=0)
                    reg_synthetic = DecisionTreeRegressor(max_depth=3).fit(X_synthetic, y_synthetic)

                    # Plot decision boundary
                    plot_step = 0.02
                    x_min, x_max = X_synthetic[:, 0].min() - 1, X_synthetic[:, 0].max() + 1
                    y_min, y_max = X_synthetic[:, 1].min() - 1, X_synthetic[:, 1].max() + 1
                    xx, yy = np.meshgrid(np.arange(x_min, x_max, plot_step),
                                        np.arange(y_min, y_max, plot_step))

                    Z = reg_synthetic.predict(np.c_[xx.ravel(), yy.ravel()])
                    Z = Z.reshape(xx.shape)

                    plt.figure(figsize=(6, 6), dpi=120)
                    plt.contourf(xx, yy, Z, cmap='coolwarm', alpha=0.8)
                    plt.scatter(X_synthetic[:, 0], X_synthetic[:, 1], c=y_synthetic, cmap='viridis', edgecolor='k', s=15)
                    plt.xlabel("Feature 1")
                    plt.ylabel("Feature 2")
                    plt.colorbar()
                    plt.show()

                    buf3 = BytesIO()
                    plt.savefig(buf3, format="png")
                    st.image(buf3)

                    # Decode categorical variables for display
                    decoded_data = X.copy()
                    for col, le in label_encoders.items():
                        if col in explanatory_vars:
                            decoded_data[col] = le.inverse_transform(X[col])

                    st.write("Decision Tree Analysis Completed in {:.2f} seconds".format(time.time() - start_time))

                    st.success("We interpret the decision tree by analyzing how the conditions split the data into branches, leading to the predicted outcomes.")

            # Spawn a new Quill editor
            st.subheader("Notes on decision tree analysis")
            dectreecontent = st_quill(placeholder="Write your notes here", key="dectreequill")

            st.session_state.dectreenotes = dectreecontent

            st.write("Decision tree analysis took {:.2f} seconds to run.".format(time.time() - start_time))

        else:

            st.warning('**Please select a methodology above to get started with the analysis.**')

        with st.expander("Build results presentation", expanded=False):

            st.info("When finished, click below to build the presentation with your results")

            buildpres = st.button("Build presentation")

            if buildpres == True:

                # Building Powerpoint presentation

                from pptx import Presentation
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.dml.color import RGBColor
                from pptx.util import Inches, Pt
                from pptx.enum.dml import MSO_THEME_COLOR
                title='   Analytics Playground\n\
                Results from analysis'
                APlogo='./Powerpoint/APlogo.png'
                ABIlogo='./Powerpoint/ABIlogo.png'
                prs = Presentation()

                # Slide 1

                # Add colour bar

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, Inches(9/1.5),Inches(16),Inches(9/8.5)
                )
                shape.shadow.inherit = False
                fill=shape.fill
                fill.solid()
                fill.fore_color.rgb=RGBColor(255,0,0)
                shape.text= title
                line=shape.line
                line.color.rgb=RGBColor(255,0,0)
                logo1=slide.shapes.add_picture(APlogo,Inches(13.5),Inches(6.0),height=Inches(1.08),width=Inches(1.0))
                logo2=slide.shapes.add_picture(ABIlogo,Inches(14.5),Inches(5.8),height=Inches(1.5),width=Inches(1.51))

                # Add text box for results

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, Inches(0.5),Inches(16),Inches(0.3))
                shape.shadow.inherit = False
                fill=shape.fill
                fill.solid()
                fill.fore_color.rgb=RGBColor(255,0,0)
                shape.text= "   Results from Exploratory Data Analysis"
                line=shape.line
                line.color.rgb=RGBColor(255,0,0)
                logo1=slide.shapes.add_picture(APlogo,Inches(14.5),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                logo2=slide.shapes.add_picture(ABIlogo,Inches(15.0),Inches(0.4),height=Inches(0.5),width=Inches(0.5))
                left = Inches(1)
                top = Inches(2)
                width = Inches(5)
                height = Inches(5)
                text_box=slide.shapes.add_textbox(left, top, width, height)
                tb=text_box.text_frame
                tb.text = st.session_state.edanotes
                prg=tb.add_paragraph()
                prg.text=" "
                prg=tb.add_paragraph()
                prg.text=''
                correlpic = slide.shapes.add_picture('correl.jpg', Inches(8), Inches(1.3), height=Inches(3.7), width=Inches(6.3))
                ppspic = slide.shapes.add_picture('pps.jpg', Inches(8), Inches(5.1), height=Inches(3.7), width=Inches(7.3))

                prs.save('EDA_presentation.pptx')

                os.startfile("EDA_presentation.pptx")

